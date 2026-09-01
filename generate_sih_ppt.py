import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Color Palette matching SIH Template & SOC Cyber Theme
    BG_COLOR = RGBColor(254, 252, 246)       # SIH Warm Cream
    TOP_BAR = RGBColor(235, 230, 218)        # Subtly darker cream
    FOOTER_BLUE = RGBColor(30, 112, 184)     # SIH Template Footer Blue
    PRIMARY_DARK = RGBColor(17, 24, 39)      # Slate 900
    TEXT_MUTED = RGBColor(75, 85, 99)        # Slate 600
    ACCENT_CYAN = RGBColor(8, 145, 178)      # Cyan 600
    ACCENT_BLUE = RGBColor(37, 99, 235)      # Blue 600
    ACCENT_GREEN = RGBColor(16, 185, 129)    # Emerald 500
    CARD_BG = RGBColor(255, 255, 255)        # Pure White Card
    CARD_BORDER = RGBColor(226, 232, 240)    # Slate 200

    def add_base_decorations(slide, slide_num, title_text, category_text="IDEA SUBMISSION"):
        # Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()

        # Top Right SIH Logo placeholder / text badge
        sih_box = slide.shapes.add_textbox(Inches(10.5), Inches(0.25), Inches(2.5), Inches(0.8))
        tf_sih = sih_box.text_frame
        tf_sih.word_wrap = True
        p_sih = tf_sih.paragraphs[0]
        p_sih.text = "SMART INDIA\nHACKATHON 2026"
        p_sih.font.bold = True
        p_sih.font.size = Pt(11)
        p_sih.font.color.rgb = PRIMARY_DARK
        p_sih.alignment = PP_ALIGN.RIGHT

        # Top Left Team Name placeholder oval
        team_oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(0.3), Inches(1.8), Inches(0.7))
        team_oval.fill.solid()
        team_oval.fill.fore_color.rgb = CARD_BG
        team_oval.line.color.rgb = RGBColor(180, 160, 200)
        team_oval.line.width = Pt(1.5)
        tf_team = team_oval.text_frame
        tf_team.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_team = tf_team.paragraphs[0]
        p_team.text = "Your Team\nName"
        p_team.font.size = Pt(9)
        p_team.font.color.rgb = PRIMARY_DARK
        p_team.alignment = PP_ALIGN.CENTER

        # Slide Main Title
        title_box = slide.shapes.add_textbox(Inches(2.7), Inches(0.3), Inches(7.6), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_t = tf_title.paragraphs[0]
        p_t.text = title_text.upper()
        p_t.font.bold = True
        p_t.font.size = Pt(22)
        p_t.font.color.rgb = PRIMARY_DARK
        p_t.alignment = PP_ALIGN.CENTER

        # Bottom Footer Bar
        footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.0), Inches(13.333), Inches(0.5))
        footer.fill.solid()
        footer.fill.fore_color.rgb = FOOTER_BLUE
        footer.line.fill.background()

        # Footer Text
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(11.0), Inches(0.4))
        p_f = footer_box.text_frame.paragraphs[0]
        p_f.text = "@SIH Idea submission - Template  |  Problem Statement: SIH26156 (NTRO Log Preprocessing)"
        p_f.font.size = Pt(10)
        p_f.font.color.rgb = RGBColor(255, 255, 255)

        # Slide Number
        num_box = slide.shapes.add_textbox(Inches(12.0), Inches(7.05), Inches(0.8), Inches(0.4))
        p_n = num_box.text_frame.paragraphs[0]
        p_n.text = str(slide_num)
        p_n.font.size = Pt(10)
        p_n.font.bold = True
        p_n.font.color.rgb = RGBColor(255, 255, 255)
        p_n.alignment = PP_ALIGN.RIGHT

    # =========================================================================
    # SLIDE 1: TITLE SLIDE
    # =========================================================================
    s1 = prs.slides.add_slide(blank_slide_layout)
    add_base_decorations(s1, 1, "IDEA TITLE: UNIVERSAL ADAPTIVE LOG PREPROCESSOR")

    # Main Banner Card
    card1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.4), Inches(10.9), Inches(5.2))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = CARD_BORDER
    card1.line.width = Pt(1.5)

    tf = card1.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.4)
    tf.margin_right = Inches(0.5)

    p0 = tf.paragraphs[0]
    p0.text = "Universal Adaptive Log Preprocessor for Heterogeneous Perimeter & Network Devices"
    p0.font.bold = True
    p0.font.size = Pt(20)
    p0.font.color.rgb = ACCENT_BLUE

    p_sub = tf.add_paragraph()
    p_sub.text = "Intelligent, 100% Air-Gapped Preprocessing Engine with Autonomous Drift Detection & Self-Healing"
    p_sub.font.size = Pt(13)
    p_sub.font.color.rgb = ACCENT_CYAN
    p_sub.space_after = Pt(14)

    bullets = [
        ("Problem Statement ID:", "SIH26156"),
        ("Problem Statement Title:", "NTRO Log Preprocessing"),
        ("Target Organization:", "National Technical Research Organisation (NTRO)"),
        ("Category / Theme:", "Security & Surveillance / Cyber Defense / Smart Automation"),
        ("Core Capability:", "Autonomous Unknown Log Detection, Adaptive Schema Inference, Zero-Cloud On-Premise Operation"),
        ("Team Details:", "Team Name: [Your Team Name]  |  Team Leader: [Name]  |  Members: [Member Names]")
    ]

    for title, val in bullets:
        p = tf.add_paragraph()
        run1 = p.add_run()
        run1.text = f"• {title} "
        run1.font.bold = True
        run1.font.size = Pt(12)
        run1.font.color.rgb = PRIMARY_DARK

        run2 = p.add_run()
        run2.text = val
        run2.font.size = Pt(12)
        run2.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(6)

    # =========================================================================
    # SLIDE 2: PROPOSED SOLUTION (Idea / Solution / Prototype)
    # =========================================================================
    s2 = prs.slides.add_slide(blank_slide_layout)
    add_base_decorations(s2, 2, "PROPOSED SOLUTION (IDEA / SOLUTION / PROTOTYPE)")

    # Left Column: Detailed Explanation & How it addresses problem
    card_l2 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.3), Inches(5.6), Inches(5.4))
    card_l2.fill.solid()
    card_l2.fill.fore_color.rgb = CARD_BG
    card_l2.line.color.rgb = CARD_BORDER
    tf_l2 = card_l2.text_frame
    tf_l2.word_wrap = True
    tf_l2.margin_left = Inches(0.3)
    tf_l2.margin_top = Inches(0.3)

    p_lh = tf_l2.paragraphs[0]
    p_lh.text = "❖ Detailed Solution & Problem Fit"
    p_lh.font.bold = True
    p_lh.font.size = Pt(15)
    p_lh.font.color.rgb = ACCENT_BLUE
    p_lh.space_after = Pt(8)

    l2_points = [
        ("The Challenge:", "Perimeter defense devices (Cisco routers, next-gen firewalls, gateways, VPNs) produce massive streams of non-standard, vendor-proprietary, and rapidly shifting logs. Hardcoded Grok/Regex parsers constantly break."),
        ("Universal Ingestion Engine:", "Accepts logs from any network device. Known formats (Cisco IOS ACLs, standard FW KV, Syslog) are recognized instantly; unseen/unknown formats are dynamically fingerprinted."),
        ("Zero-Code Adaptive Parsing:", "Extracts IPs, ports, action verbs (ALLOW/DENY/DROP), protocols, and timestamps. Uses semantic synonym inference (SRCIP -> source_ip, DP -> destination_port) without writing custom rules."),
        ("Universal Normalization:", "Converts all heterogeneous inputs into a standardized JSON defense schema for downstream SIEM, SOC, and anomaly detection.")
    ]
    for h, b in l2_points:
        p = tf_l2.add_paragraph()
        r1 = p.add_run()
        r1.text = f"• {h} "
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = PRIMARY_DARK
        r2 = p.add_run()
        r2.text = b
        r2.font.size = Pt(10)
        r2.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(5)

    # Right Column: Innovation & Uniqueness
    card_r2 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.3), Inches(5.7), Inches(5.4))
    card_r2.fill.solid()
    card_r2.fill.fore_color.rgb = CARD_BG
    card_r2.line.color.rgb = CARD_BORDER
    tf_r2 = card_r2.text_frame
    tf_r2.word_wrap = True
    tf_r2.margin_left = Inches(0.3)
    tf_r2.margin_top = Inches(0.3)

    p_rh = tf_r2.paragraphs[0]
    p_rh.text = "❖ Innovation & Key R&D Differentiators"
    p_rh.font.bold = True
    p_rh.font.size = Pt(15)
    p_rh.font.color.rgb = ACCENT_CYAN
    p_rh.space_after = Pt(8)

    r2_points = [
        ("1. Autonomous Self-Healing:", "When firmware updates mutate log keys (e.g. ACTION= -> ACT=, SRC= -> SRCIP=), the Drift Detector flags the confidence drop (100% -> 25%) and dynamically synthesizes Parser v2.0 to restore 100% confidence."),
        ("2. Deterministic Confidence Scorer:", "Explainable multi-attribute weighting (Timestamp +20%, Src IP +20%, Dst IP +20%, Protocol +15%, Action +15%, Port +10%) eliminating black-box stochastic guessing."),
        ("3. 100% Air-Gapped Local Processing:", "Zero cloud dependency, zero external LLM API calls. Sensitive defense network topology and internal IPs remain completely secure on-premises."),
        ("4. Explainable SOC Visualizer:", "Interactive React/Tailwind dashboard displaying real-time pipeline topology, semantic mapping trees, drift timelines, and telemetry analytics.")
    ]
    for h, b in r2_points:
        p = tf_r2.add_paragraph()
        r1 = p.add_run()
        r1.text = f"★ {h} "
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = PRIMARY_DARK
        r2 = p.add_run()
        r2.text = b
        r2.font.size = Pt(10)
        r2.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(5)

    # =========================================================================
    # SLIDE 3: TECHNICAL ARCHITECTURE & PIPELINE
    # =========================================================================
    s3 = prs.slides.add_slide(blank_slide_layout)
    add_base_decorations(s3, 3, "TECHNICAL ARCHITECTURE & PROCESSING PIPELINE")

    # Pipeline steps horizontal boxes
    steps = [
        ("1. Ingest", "Local Raw Stream", ACCENT_BLUE),
        ("2. Fingerprint", "Structural Delimiters", ACCENT_CYAN),
        ("3. Classify", "Known vs Unknown", RGBColor(99, 102, 241)),
        ("4. Adaptive Parse", "Semantic Synonym Map", RGBColor(139, 92, 246)),
        ("5. Confidence", "Explainable Scorer", ACCENT_GREEN),
        ("6. Normalize", "Universal Defense Event", ACCENT_GREEN),
        ("7. Self-Heal", "Drift Anomaly v1 -> v2", RGBColor(245, 158, 11))
    ]

    box_w = Inches(1.55)
    box_h = Inches(1.2)
    start_x = Inches(0.8)
    gap = Inches(0.15)

    for i, (title, sub, col) in enumerate(steps):
        x = start_x + i * (box_w + gap)
        bx = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.3), box_w, box_h)
        bx.fill.solid()
        bx.fill.fore_color.rgb = CARD_BG
        bx.line.color.rgb = col
        bx.line.width = Pt(2)

        tf_b = bx.text_frame
        tf_b.word_wrap = True
        tf_b.margin_top = Inches(0.15)
        p1 = tf_b.paragraphs[0]
        p1.text = title
        p1.font.bold = True
        p1.font.size = Pt(11)
        p1.font.color.rgb = col
        p1.alignment = PP_ALIGN.CENTER

        p2 = tf_b.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(9)
        p2.font.color.rgb = TEXT_MUTED
        p2.alignment = PP_ALIGN.CENTER

    # Lower Left Card: Tech Stack
    card_tech = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.7), Inches(5.6), Inches(4.0))
    card_tech.fill.solid()
    card_tech.fill.fore_color.rgb = CARD_BG
    card_tech.line.color.rgb = CARD_BORDER
    tf_tech = card_tech.text_frame
    tf_tech.word_wrap = True
    tf_tech.margin_left = Inches(0.3)
    tf_tech.margin_top = Inches(0.25)

    p_th = tf_tech.paragraphs[0]
    p_th.text = "❖ Core Technology Stack"
    p_th.font.bold = True
    p_th.font.size = Pt(14)
    p_th.font.color.rgb = ACCENT_BLUE
    p_th.space_after = Pt(6)

    stack = [
        ("Backend Core:", "Python 3.13, FastAPI (High-performance asynchronous REST API), Pydantic v2 (Strict Schema Validation)"),
        ("Parsing & Inference:", "Pattern-based Feature Extractor, Semantic Synonym Dictionaries, Positional Token Heuristics, Dynamic AST Generator"),
        ("Frontend Visualizer:", "React 18, Vite, Tailwind CSS (Dark SOC Theme), Recharts (Real-time Telemetry Charts), Lucide Icons"),
        ("Air-Gap Security:", "100% On-Premise Execution, Zero External Network Egress, In-Memory & Local SQLite Store")
    ]
    for k, v in stack:
        p = tf_tech.add_paragraph()
        r1 = p.add_run()
        r1.text = f"• {k} "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = PRIMARY_DARK
        r2 = p.add_run()
        r2.text = v
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(4)

    # Lower Right Card: Universal Schema Output
    card_sch = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.7), Inches(5.7), Inches(4.0))
    card_sch.fill.solid()
    card_sch.fill.fore_color.rgb = CARD_BG
    card_sch.line.color.rgb = CARD_BORDER
    tf_sch = card_sch.text_frame
    tf_sch.word_wrap = True
    tf_sch.margin_left = Inches(0.3)
    tf_sch.margin_top = Inches(0.25)

    p_sh = tf_sch.paragraphs[0]
    p_sh.text = "❖ Universal Normalized Defense Event Schema"
    p_sh.font.bold = True
    p_sh.font.size = Pt(14)
    p_sh.font.color.rgb = ACCENT_CYAN
    p_sh.space_after = Pt(6)

    sch_text = """{\n  "timestamp": "2026-08-31T21:13:04Z",\n  "device": "EDGE-X / FW01 / Cisco-Router",\n  "device_type": "Perimeter Gateway",\n  "source_ip": "172.20.1.50",\n  "destination_ip": "10.10.4.8",\n  "source_port": 52144,\n  "destination_port": 3389,\n  "protocol": "TCP",\n  "action": "DENY",\n  "parser_type": "adaptive | known | self_healed",\n  "parser_version": "fw_v2.0 (Self-Healed)",\n  "confidence": 100.0,\n  "status": "self_healed"\n}"""

    p_code = tf_sch.add_paragraph()
    p_code.text = sch_text
    p_code.font.size = Pt(8.5)
    p_code.font.name = "Courier New"
    p_code.font.color.rgb = RGBColor(16, 185, 129)

    # =========================================================================
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # =========================================================================
    s4 = prs.slides.add_slide(blank_slide_layout)
    add_base_decorations(s4, 4, "FEASIBILITY AND VIABILITY")

    cards_data_4 = [
        ("Analysis of Feasibility", ACCENT_BLUE, [
            ("Technical Feasibility:", "Lightweight regex compilation and deterministic semantic lookup achieve sub-millisecond per-log processing latency."),
            ("Operational Feasibility:", "Drop-in compatibility with existing Syslog daemons (rsyslog, syslog-ng), Logstash, FluentBit, and Kafka streams."),
            ("Economic Viability:", "Zero recurring cloud LLM API costs or token fees. Completely free from proprietary vendor licensing.")
        ]),
        ("Potential Challenges & Risks", RGBColor(239, 68, 68), [
            ("High EPS Load (100k+ Logs/sec):", "Risk of CPU saturation during continuous real-time semantic inference."),
            ("Severely Obfuscated Formats:", "Risk of missing field semantics in completely non-standard proprietary binaries or compressed blobs."),
            ("False Positive Drift Triggers:", "Minor intermittent log corruption might trigger unnecessary parser re-synthesis.")
        ]),
        ("Mitigation & Engineering Strategies", ACCENT_GREEN, [
            ("Vectorized Regex & Hyperscan:", "Compile parser ASTs into hardware-accelerated C-bindings (Intel Hyperscan) for 100,000+ EPS throughput."),
            ("Positional Heuristic Fallback:", "Dual-tier validation combines key-value matching with positional token entropy and IP/Port regexes."),
            ("Confidence Thresholding & Quorum:", "Require minimum batch verification (e.g. 5+ failed logs) before promoting self-healed Parser v2.0.")
        ])
    ]

    col_w = Inches(3.7)
    col_gap = Inches(0.3)
    start_x_4 = Inches(0.8)

    for i, (head, col, items) in enumerate(cards_data_4):
        x = start_x_4 + i * (col_w + col_gap)
        c = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.3), col_w, Inches(5.4))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = CARD_BORDER
        tf_c = c.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = Inches(0.25)
        tf_c.margin_top = Inches(0.25)
        tf_c.margin_right = Inches(0.25)

        ph = tf_c.paragraphs[0]
        ph.text = f"❖ {head}"
        ph.font.bold = True
        ph.font.size = Pt(13)
        ph.font.color.rgb = col
        ph.space_after = Pt(8)

        for t, d in items:
            p = tf_c.add_paragraph()
            r1 = p.add_run()
            r1.text = f"• {t} "
            r1.font.bold = True
            r1.font.size = Pt(10)
            r1.font.color.rgb = PRIMARY_DARK
            r2 = p.add_run()
            r2.text = d
            r2.font.size = Pt(9.5)
            r2.font.color.rgb = TEXT_MUTED
            p.space_after = Pt(6)

    # =========================================================================
    # SLIDE 5: IMPACT, BENEFITS, AND POTENTIAL APPLICATIONS
    # =========================================================================
    s5 = prs.slides.add_slide(blank_slide_layout)
    add_base_decorations(s5, 5, "IMPACT, BENEFITS, AND POTENTIAL APPLICATIONS")

    # 4 Quadrant Cards
    quads = [
        ("National Security & Defense (NTRO / CERT-In / Military Grid)", ACCENT_BLUE, [
            "Seamlessly integrates heterogeneous perimeter hardware across border routers, naval gateways, and military SOCs.",
            "100% Air-gapped deployment prevents sensitive defense network topology and firewall configurations from leaking to external clouds."
        ]),
        ("Critical Infrastructure & Financial Institutions", ACCENT_CYAN, [
            "Protects Banking networks, Power Grids (SCADA/ICS), and Telecom backbones from log pipeline blackouts during vendor firmware updates.",
            "Normalizes multi-vendor perimeter traffic (Palo Alto, Fortinet, Cisco, Check Point) into a single unified format."
        ]),
        ("Massive Operational Efficiency & Cost Reduction", ACCENT_GREEN, [
            "Reduces manual Grok/Regex parser engineering time by over 90%.",
            "Eliminates multi-million rupee cloud LLM token expenses through deterministic local edge processing."
        ]),
        ("Resilient Self-Healing Zero-Downtime Pipeline", RGBColor(245, 158, 11), [
            "Eliminates silent data corruption caused by schema drift.",
            "Automated transition from Parser v1 to v2 ensures continuous uninterrupted log visibility for threat hunting."
        ])
    ]

    q_w = Inches(5.6)
    q_h = Inches(2.55)
    positions = [
        (Inches(0.8), Inches(1.3)),
        (Inches(6.8), Inches(1.3)),
        (Inches(0.8), Inches(4.1)),
        (Inches(6.8), Inches(4.1))
    ]

    for (head, col, pts), (x, y) in zip(quads, positions):
        c = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, q_w, q_h)
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = CARD_BORDER
        tf_q = c.text_frame
        tf_q.word_wrap = True
        tf_q.margin_left = Inches(0.3)
        tf_q.margin_top = Inches(0.2)
        tf_q.margin_right = Inches(0.3)

        ph = tf_q.paragraphs[0]
        ph.text = f"❖ {head}"
        ph.font.bold = True
        ph.font.size = Pt(12)
        ph.font.color.rgb = col
        ph.space_after = Pt(4)

        for pt in pts:
            p = tf_q.add_paragraph()
            p.text = f"• {pt}"
            p.font.size = Pt(9.5)
            p.font.color.rgb = TEXT_MUTED
            p.space_after = Pt(3)

    # =========================================================================
    # SLIDE 6: RESEARCH AND REFERENCES
    # =========================================================================
    s6 = prs.slides.add_slide(blank_slide_layout)
    add_base_decorations(s6, 6, "RESEARCH AND REFERENCES")

    card_ref = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.3), Inches(11.3), Inches(5.4))
    card_ref.fill.solid()
    card_ref.fill.fore_color.rgb = CARD_BG
    card_ref.line.color.rgb = CARD_BORDER
    tf_ref = card_ref.text_frame
    tf_ref.word_wrap = True
    tf_ref.margin_left = Inches(0.4)
    tf_ref.margin_top = Inches(0.3)

    p_rh = tf_ref.paragraphs[0]
    p_rh.text = "❖ Technical References & Research Foundation"
    p_rh.font.bold = True
    p_rh.font.size = Pt(15)
    p_rh.font.color.rgb = ACCENT_BLUE
    p_rh.space_after = Pt(10)

    refs = [
        ("IETF RFC 5424 & RFC 3164:", "The Syslog Protocol & BSD Syslog Standards for Network Device Logging."),
        ("NIST Special Publication 800-92:", "Guide to Computer Security Log Management - Recommendations of the National Institute of Standards and Technology."),
        ("Drain Log Parsing Research (He et al., IEEE ICWS):", "Drain: An Online Log Parsing Approach with Fixed Depth Trees for High-Volume Telemetry Systems."),
        ("Spell & Logram Mining (Du et al., IEEE ICDM / Dai et al., IEEE TSE):", "Streaming Parsing of Unstructured System Logs and Log-based Automated Anomaly Detection."),
        ("Schema Drift & Data Quality (Gao et al., ACM SIGMOD):", "Adaptive Schema Evolution and Anomaly Detection in Heterogeneous Streaming Pipelines."),
        ("Smart India Hackathon 2026 Problem Guidelines:", "Problem Statement SIH26156: 'NTRO Log Preprocessing' - National Technical Research Organisation.")
    ]

    for title, desc in refs:
        p = tf_ref.add_paragraph()
        r1 = p.add_run()
        r1.text = f"• {title} "
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = PRIMARY_DARK
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(8)

    output_path = r"c:\Users\gauta\Log Prototype\SIH26156_NTRO_Log_Preprocessor_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_presentation()
